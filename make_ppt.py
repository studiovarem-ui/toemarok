#!/usr/bin/env python3
"""퇴마록 재현용 프롬프트 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG_DARK = RGBColor(26, 10, 10)       # #1a0a0a
BG_CARD = RGBColor(42, 28, 20)       # #2a1c14
GOLD = RGBColor(255, 215, 0)         # #FFD700
RED = RGBColor(221, 68, 68)          # #DD4444
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
MED_GRAY = RGBColor(150, 150, 150)
DIM_GRAY = RGBColor(100, 100, 100)
GREEN = RGBColor(68, 221, 68)
BLUE = RGBColor(100, 180, 255)
ORANGE = RGBColor(255, 140, 40)
PURPLE = RGBColor(170, 100, 255)
CYAN = RGBColor(130, 220, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = "맑은 고딕"
    return txBox


def add_multiline(slide, left, top, width, height, lines, default_size=13, default_color=LIGHT_GRAY):
    """lines: list of (text, size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "맑은 고딕"
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def add_table(slide, left, top, width, height, rows, col_widths=None):
    """rows: list of lists. First row = header."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.name = "맑은 고딕"
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = GOLD
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.alignment = PP_ALIGN.CENTER

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r_idx == 0:
                cell_fill.fore_color.rgb = RGBColor(50, 30, 20)
            elif r_idx % 2 == 0:
                cell_fill.fore_color.rgb = RGBColor(35, 25, 18)
            else:
                cell_fill.fore_color.rgb = RGBColor(45, 30, 22)

    return table_shape


# ════════════════════════════════════════════
# SLIDE 1: 표지
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Decorative top bar
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

# Title
add_text(slide, Inches(0), Inches(1.8), SW, Inches(1.2),
         "퇴  마  록", 72, GOLD, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.0), SW, Inches(0.5),
         "TOEMAROK", 24, RGBColor(136, 102, 68), False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.5), SW, Inches(0.5),
         "한국 신화 뱀서라이크 — 재현용 프롬프트 가이드", 16, MED_GRAY, False, PP_ALIGN.CENTER)

# Bottom info
add_text(slide, Inches(0), Inches(5.8), SW, Inches(0.4),
         "HTML5 Canvas  •  Pixel Art  •  No External Assets  •  Mobile Optimized", 12, DIM_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(6.2), SW, Inches(0.4),
         "github.com/studiovarem-ui/toemarok", 11, RGBColor(100, 140, 180), False, PP_ALIGN.CENTER)

# Bottom bar
add_shape(slide, Inches(0), Inches(7.44), SW, Inches(0.06), GOLD)


# ════════════════════════════════════════════
# SLIDE 2: 기술 스펙 개요
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "📋 기술 스펙 & 핵심 구조", 28, GOLD, True)

# Left column
left_lines = [
    ("▎ 캔버스 & 렌더링", 14, ORANGE, True),
    ("  • 해상도: 400×700 (세로형 모바일)", 12, LIGHT_GRAY),
    ("  • CSS flex 중앙정렬, image-rendering: pixelated", 12, LIGHT_GRAY),
    ("  • 모든 그래픽 Canvas API 직접 렌더링 (외부 이미지 없음)", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 게임 상태 머신", 14, ORANGE, True),
    ("  title → charSelect → playing ↔ levelUp → gameOver / victory", 12, CYAN),
    ("", 8),
    ("▎ 오디오", 14, ORANGE, True),
    ("  • Web Audio API 프로시저럴 사운드", 12, LIGHT_GRAY),
    ("  • 종류: hit, kill, levelup, bomb, pickup, boss", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 저장 시스템", 14, ORANGE, True),
    ("  • localStorage: unlocks, bestTime, bestKills, totalClears", 12, LIGHT_GRAY),
]
add_multiline(slide, Inches(0.5), Inches(1.0), Inches(5.8), Inches(5.5), left_lines)

# Right column
right_lines = [
    ("▎ 모바일 입력 (중요!)", 14, RED, True),
    ("  • touch-area div (z-index:10) 위에서 터치 이벤트 처리", 12, LIGHT_GRAY),
    ("  • screenToCanvas(): canvas.getBoundingClientRect() 사용", 12, CYAN),
    ("  • touchstart에서 e.preventDefault() → click 차단됨", 12, LIGHT_GRAY),
    ("  • handleTap() 함수로 touch/click 양쪽에서 통합 호출", 12, GREEN),
    ("", 8),
    ("▎ 핵심 버그 방지 패턴", 14, RED, True),
    ("  • 적 고유 uid (enemyIdCounter++) — 배열 index 사용 금지", 12, LIGHT_GRAY),
    ("  • 존/궤도 무기의 _tick 추적은 uid 기반", 12, LIGHT_GRAY),
    ("  • 천둥 슬로우: setTimeout ❌ → 게임타임 타이머 ✅", 12, LIGHT_GRAY),
    ("  • 호밍 투사체: 타겟 사망시 자동 리타겟", 12, LIGHT_GRAY),
    ("  • 적 존 대미지: dmg*dt ❌ → 0.5초 틱 ✅", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 파일 구조", 14, ORANGE, True),
    ("  index.html  /  style.css  /  game.js (~2100줄)", 12, LIGHT_GRAY),
]
add_multiline(slide, Inches(6.8), Inches(1.0), Inches(6.0), Inches(5.5), right_lines)


# ════════════════════════════════════════════
# SLIDE 3: 캐릭터 6종
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "🎭 캐릭터 6종", 28, GOLD, True)

chars = [
    ["이름", "설명", "시작무기", "HP", "SPD", "ATK", "RANGE", "해금조건"],
    ["퇴마사", "균형형", "부적(0)", "150", "130", "1.2", "50", "기본"],
    ["무녀", "원거리", "신령방울(1)", "120", "120", "1.3", "60", "기본"],
    ["도깨비", "근접", "방망이(2)", "180", "110", "1.5", "40", "1클리어"],
    ["구미호", "속도", "여우불(3)", "100", "170", "1.1", "45", "3클리어"],
    ["장군", "탱커", "신궁(5)", "220", "100", "1.2", "35", "5클리어"],
    ["산신령", "소환", "천둥(4)", "130", "110", "1.4", "70", "생존900초"],
]
add_table(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.2), chars)

# Sprite descriptions
add_text(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
         "스프라이트 디자인 (16-20px 픽셀아트)", 16, ORANGE, True)

sprite_lines = [
    ("퇴마사   하얀 도복, 검은 갓모자, 파란 허리띠, 손에 빨간 부적", 11, LIGHT_GRAY),
    ("무녀      빨간 치마 흰 저고리, 긴 검은 머리+분홍 꽃장식, 손에 금색 방울", 11, LIGHT_GRAY),
    ("도깨비   파란 피부, 금색 뿔, 빨간 눈, 큰 갈색 방망이 (쇠스파이크)", 11, LIGHT_GRAY),
    ("구미호   분홍 한복, 여우귀, 5개 주황꼬리 (흔들림 애니메이션), 여우불 이펙트", 11, LIGHT_GRAY),
    ("장군      갈색 갑옷+금장식, 빨간 투구, 창", 11, LIGHT_GRAY),
    ("산신령   흰 도복, 긴 흰 수염, 녹색 오라, 옆에 주황 호랑이", 11, LIGHT_GRAY),
]
add_multiline(slide, Inches(0.7), Inches(4.0), Inches(11.5), Inches(3.0), sprite_lines)


# ════════════════════════════════════════════
# SLIDE 4: 무기 8종
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "⚔️ 무기 8종", 28, GOLD, True)

weapons = [
    ["ID", "이름", "타입", "쿨타임", "설명", "진화"],
    ["0", "부적", "projectile", "0.5", "전방 투사체, 레벨=개수", "봉인진: 처치시 범위폭발"],
    ["1", "신령 방울", "homing", "0.8", "유도탄 min(lv,3)발", "—"],
    ["2", "도깨비 방망이", "spin", "1.0", "궤도 회전, min(1+lv/2,3)개", "여의봉: 3배크기"],
    ["3", "여우불", "zone", "1.5", "불꽃 장판, min(1+lv/2,3)개", "삼매화: 9개궤도화염"],
    ["4", "천둥", "thunder", "1.2", "즉발 낙뢰 min(lv,3)대상", "뇌신: 5체인+마비"],
    ["5", "신궁", "pierce", "0.6", "관통 화살, lv5=무한관통", "—"],
    ["6", "용의 숨결", "breath", "1.8", "전방 부채꼴, lv3+화상", "청룡: 화면관통용"],
    ["7", "귀살검", "slash", "0.4", "전방 베기, lv5:후방도 공격", "—"],
]
add_table(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(2.8), weapons)

# Projectile rendering
add_text(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
         "투사체 렌더링 스타일", 16, ORANGE, True)

proj_lines = [
    ("부적       회전하는 빨간 카드 + 금색 부적문양 + 잔상 트레일", 11, LIGHT_GRAY),
    ("방울       금색 구슬 + 음파 링 이펙트 (확장되는 원)", 11, LIGHT_GRAY),
    ("방망이    궤도 회전 곤봉 (나무 손잡이 + 쇠머리 + 스파이크) + 궤도링 표시", 11, LIGHT_GRAY),
    ("여우불    다층 불꽃 (#FF4400→#FF8844→#FFCC44→#FFF) + 불씨 파티클", 11, LIGHT_GRAY),
    ("천둥       지그재그 번개 볼트 (fillPath) + 글로우 (반투명 원)", 11, LIGHT_GRAY),
    ("화살       화살촉 삼각형 + 속도선 잔상 (반투명 가로줄)", 11, LIGHT_GRAY),
    ("용숨결    evolved: 큰 파란 드래곤 투사체 (3원 몸통 + 노란 눈)", 11, LIGHT_GRAY),
    ("귀살검    반달형 슬래시 아크 (ctx.arc stroke) + 흰색 하이라이트", 11, LIGHT_GRAY),
]
add_multiline(slide, Inches(0.7), Inches(4.5), Inches(11.5), Inches(3.0), proj_lines)


# ════════════════════════════════════════════
# SLIDE 5: 패시브 & 진화
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "💎 패시브 8종 & 진화 5종", 28, GOLD, True)

passives = [
    ["ID", "이름", "스탯", "레벨당", "진화 조합"],
    ["0", "음양오행", "atk", "+15%", "부적 Lv5 → 봉인진"],
    ["1", "구미호 가죽", "spd", "+10%", "여우불 Lv5 → 삼매화"],
    ["2", "풍백 가호", "range", "+12%", "천둥 Lv5 → 뇌신"],
    ["3", "황금", "exp", "+15%", "방망이 Lv5 → 여의봉"],
    ["4", "산삼", "regen", "3초/lv, 3+lv*2 회복", "—"],
    ["5", "여의주", "cdr", "-8%", "용숨결 Lv5 → 청룡"],
    ["6", "도깨비감투", "dodge", "+10%", "—"],
    ["7", "삼족오 부적", "crit", "+8%", "—"],
]
add_table(slide, Inches(0.3), Inches(1.0), Inches(7.5), Inches(2.8), passives)

# Visual effects
effect_lines = [
    ("▎ 패시브 시각 효과 (중요!)", 14, ORANGE, True),
    ("", 6),
    ("  산삼 회복     초록색 \"+N\" 숫자 (#44FF44) + 초록 파티클", 12, GREEN),
    ("  회피 발동     \"회피!\" 하늘색 텍스트 (#88FFFF) + 파란 파티클", 12, CYAN),
    ("  크리티컬      \"크리!\" 노란 라벨 (#FFDD00) + 큰 대미지 숫자", 12, GOLD),
    ("", 8),
    ("▎ 크리티컬 시스템", 14, ORANGE, True),
    ("  • rollCrit() 함수로 투사체별 개별 판정", 12, LIGHT_GRAY),
    ("  • 크리시 대미지 ×2", 12, LIGHT_GRAY),
]
add_multiline(slide, Inches(8.3), Inches(1.0), Inches(4.5), Inches(3.5), effect_lines)

# Evolutions
evo_lines = [
    ("▎ 진화 조건: 무기 Lv5 + 해당 패시브 보유", 14, RED, True),
    ("", 6),
    ("  봉인진   부적(0) + 음양오행(0) → 7발 + 처치시 범위폭발", 12, PURPLE),
    ("  삼매화   여우불(3) + 구미호가죽(1) → 9개 궤도 화염", 12, PURPLE),
    ("  뇌신      천둥(4) + 풍백가호(2) → 5체인 + 슬로우(1.5초)", 12, PURPLE),
    ("  여의봉   방망이(2) + 황금(3) → 4궤도 + 넉백", 12, PURPLE),
    ("  청룡      용숨결(6) + 여의주(5) → 화면 관통 드래곤", 12, PURPLE),
]
add_multiline(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(3.0), evo_lines)


# ════════════════════════════════════════════
# SLIDE 6: 적 8종
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "👹 적 8종 + 보스 2종", 28, GOLD, True)

enemies = [
    ["이름", "HP", "SPD", "DMG", "패턴", "등장", "스프라이트"],
    ["잡귀", "2", "65", "2", "straight", "0초", "보라 유령, 빨간 눈"],
    ["도깨불", "4", "50", "4", "zigzag", "30초", "파란 도깨비불"],
    ["물귀신", "6", "40", "4", "aimed", "90초", "녹색 물괴물, 해초머리"],
    ["야차", "5", "100", "5", "swooper", "150초", "빨간 악귀, 금색 뿔"],
    ["강시", "16", "30", "6", "tank", "210초", "청록 관복, 노란 부적"],
    ["원귀", "6", "18", "4", "sniper(투사체)", "270초", "반투명 흰 유령"],
    ["삼두구", "3", "55", "3", "formation(3동시)", "330초", "3머리 개"],
    ["이무기", "12", "45", "5", "spiral", "400초", "녹색 뱀용, 뿔"],
]
add_table(slide, Inches(0.3), Inches(1.0), Inches(12.7), Inches(2.8), enemies)

# Bosses
bosses = [
    ["이름", "HP", "SPD", "DMG", "등장", "패턴"],
    ["귀왕", "200", "35", "12", "450초", "추적→돌진→범위폭발+잡귀3소환"],
    ["구미호왕", "500", "30", "15", "900초", "궤도+유도탄→여우불5장판"],
]
add_table(slide, Inches(0.3), Inches(4.1), Inches(10), Inches(1.1), bosses)

boss_lines = [
    ("▎ 보스 스프라이트", 14, ORANGE, True),
    ("  귀왕: 2~3배 크기, 5개 금색 뿔 왕관, 빛나는 주황 눈, 큰 도끼", 12, LIGHT_GRAY),
    ("  구미호왕: 거대 9미호, 9개 불꽃 꼬리, 금색 한복, 여우귀", 12, LIGHT_GRAY),
    ("", 6),
    ("▎ 적 공통 시스템", 14, ORANGE, True),
    ("  • 피격 시 흰색 플래시 (hitFlash 0.15초)", 12, LIGHT_GRAY),
    ("  • 대미지 입은 적만 HP바 표시 (16×2px)", 12, LIGHT_GRAY),
    ("  • 사망 시 파티클 폭발 + EXP 오브 드롭", 12, LIGHT_GRAY),
    ("  • 구미호왕 처치 → victory 상태 (승리화면)", 12, GREEN),
]
add_multiline(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(2.5), boss_lines)


# ════════════════════════════════════════════
# SLIDE 7: 밸런스 & 스폰
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "⚖️ 밸런스 & 스폰 시스템", 28, GOLD, True)

bal_left = [
    ("▎ 스폰 레이트", 16, ORANGE, True),
    ("", 4),
    ("  baseRate = 1.2 + gameTime/60 * 0.6", 13, CYAN),
    ("           + floor(gameTime/60) * 0.3", 13, CYAN),
    ("", 6),
    ("  → 시작: ~1.2마리/초", 12, LIGHT_GRAY),
    ("  → 5분: ~2.7마리/초", 12, LIGHT_GRAY),
    ("  → 10분: ~4.2마리/초", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 웨이브 버스트 (30초마다)", 16, ORANGE, True),
    ("", 4),
    ("  burstCount = 4 + floor(gameTime/60)", 13, CYAN),
    ("  → 잡귀를 원형 배치로 한꺼번에 소환", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ HP 스케일링", 16, ORANGE, True),
    ("", 4),
    ("  hpScale = 1 + gameTime/60 * 0.1", 13, CYAN),
    ("  → 10분에 적 HP +100%", 12, LIGHT_GRAY),
]
add_multiline(slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(5.5), bal_left)

bal_right = [
    ("▎ 맵", 16, ORANGE, True),
    ("", 4),
    ("  • 반경 1200 원형 맵", 12, LIGHT_GRAY),
    ("  • 40×40 타일, 해시 기반 색상 변화 + 풀/돌 장식", 12, LIGHT_GRAY),
    ("  • 경계 접근시 빨간 원 표시", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 적 제한", 16, ORANGE, True),
    ("", 4),
    ("  • 최대 200마리", 12, LIGHT_GRAY),
    ("  • 초과시 가장 먼 적 제거 (sort+pop)", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 경험치", 16, ORANGE, True),
    ("", 4),
    ("  • 적 사망 → EXP 오브 드롭", 12, LIGHT_GRAY),
    ("  • 자석 범위 = range × rangeMul", 12, LIGHT_GRAY),
    ("  • expToNext = level*10 + level^1.5 * 5", 12, CYAN),
    ("", 8),
    ("▎ 폭탄", 16, ORANGE, True),
    ("", 4),
    ("  • 쿨다운 30초, 전체 적 50대미지 (보스 25)", 12, LIGHT_GRAY),
    ("  • 금색 화면 플래시 + 파티클 30개", 12, LIGHT_GRAY),
]
add_multiline(slide, Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.8), bal_right)


# ════════════════════════════════════════════
# SLIDE 8: UI & 컨트롤
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "🎮 UI & 모바일 컨트롤", 28, GOLD, True)

ui_left = [
    ("▎ HUD 레이아웃", 16, ORANGE, True),
    ("", 4),
    ("  좌상단: 타이머 (MM:SS, 금색, 14px)", 12, LIGHT_GRAY),
    ("  우상단: 킬 수 (💀 N, 빨간색)", 12, LIGHT_GRAY),
    ("  중앙 상단: 현재 적 수 (N 요괴, 회색)", 12, LIGHT_GRAY),
    ("  하단 좌: HP바 (100px, 초록/빨강) + Lv 표시", 12, LIGHT_GRAY),
    ("  하단: EXP바 (전체 너비, 민트색)", 12, LIGHT_GRAY),
    ("  좌하단: 무기 아이콘 (22×22 사각형, 이름+레벨)", 12, LIGHT_GRAY),
    ("  우하단: 폭탄 버튼 (원형, 반경 28px)", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 화면 플래시 이펙트", 16, ORANGE, True),
    ("", 4),
    ("  피격: 빨간색 (0.15초)", 12, RGBColor(255, 100, 100)),
    ("  폭탄: 금색 (0.5초)", 12, GOLD),
    ("  보스 처치: 금색 (0.3초)", 12, GOLD),
    ("  승리: 금색 (1.0초)", 12, GOLD),
]
add_multiline(slide, Inches(0.5), Inches(1.0), Inches(5.8), Inches(5.5), ui_left)

ui_right = [
    ("▎ 모바일 컨트롤", 16, ORANGE, True),
    ("", 4),
    ("  이동: 터치 드래그 (조이스틱식, 5px 데드존)", 12, LIGHT_GRAY),
    ("  폭탄: 우하단 원형 버튼 터치", 12, LIGHT_GRAY),
    ("  메뉴: handleTap()으로 터치/클릭 통합", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 키보드 컨트롤", 16, ORANGE, True),
    ("", 4),
    ("  이동: 방향키 / WASD", 12, LIGHT_GRAY),
    ("  폭탄: 스페이스바", 12, LIGHT_GRAY),
    ("  레벨업: 1, 2, 3 키", 12, LIGHT_GRAY),
    ("  시작: Enter / Space", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 폭탄 버튼 렌더링", 16, ORANGE, True),
    ("", 4),
    ("  대기: 금색 글로우 + '부' 글자 + '폭탄' 라벨", 12, LIGHT_GRAY),
    ("  쿨다운: 회색 원 + 쿨다운 아크 + 잔여시간 숫자", 12, LIGHT_GRAY),
    ("", 8),
    ("▎ 레벨업 UI", 16, ORANGE, True),
    ("", 4),
    ("  반투명 검정 오버레이 + 3장 카드 (110×160px)", 12, LIGHT_GRAY),
    ("  진화 카드: 보라색 배경, 금색 텍스트", 12, PURPLE),
]
add_multiline(slide, Inches(7.0), Inches(1.0), Inches(5.8), Inches(6.0), ui_right)


# ════════════════════════════════════════════
# SLIDE 9: 전체 프롬프트 (1/2)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), RED)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "📝 재현용 프롬프트 (1/2)", 28, RED, True)

prompt1 = """한국 신화 뱀서라이크(Vampire Survivors류) 게임 "퇴마록"을 HTML5 Canvas 단일 파일(game.js)로 만들어줘.
외부 이미지/라이브러리 없이 Canvas API로 모든 픽셀아트를 직접 그려야 해.
GitHub Pages로 배포할 거야. repo: studiovarem-ui/toemarok

## 기술 스펙
- 캔버스: 400×700, CSS flex 중앙정렬, image-rendering: pixelated
- 모바일 최적화: touch-area div(z-index:10), canvas.getBoundingClientRect()로 좌표 변환
- touchstart에서 e.preventDefault() 사용하므로 click 이벤트 차단됨 → handleTap() 함수를 만들어서 touchstart와 click 양쪽에서 호출
- 상태: title → charSelect → playing ↔ levelUp → gameOver / victory
- Web Audio API로 효과음 (hit, kill, levelup, bomb, pickup, boss)
- localStorage 저장 (unlocks, bestTime, bestKills, totalClears)

## 캐릭터 6종 (각각 고유 픽셀아트 스프라이트 함수)
퇴마사(균형, 부적(0), HP:150 SPD:130 ATK:1.2 RANGE:50, 기본해금)
무녀(원거리, 신령방울(1), HP:120 SPD:120 ATK:1.3 RANGE:60, 기본해금)
도깨비(근접, 방망이(2), HP:180 SPD:110 ATK:1.5 RANGE:40, 1클리어)
구미호(속도, 여우불(3), HP:100 SPD:170 ATK:1.1 RANGE:45, 3클리어)
장군(탱커, 신궁(5), HP:220 SPD:100 ATK:1.2 RANGE:35, 5클리어)
산신령(소환, 천둥(4), HP:130 SPD:110 ATK:1.4 RANGE:70, 생존900초)

각 캐릭터는 16-20px 상세 픽셀아트: 퇴마사(하얀 도복+검은 갓+부적), 무녀(빨간치마+꽃장식+방울),
도깨비(파란피부+금뿔+방망이), 구미호(분홍한복+5꼬리애니+여우불),
장군(갈색갑옷+빨간투구+창), 산신령(흰도복+흰수염+녹색오라+호랑이)"""

add_multiline(slide, Inches(0.4), Inches(0.9), Inches(12.5), Inches(6.3), [
    (prompt1, 10, LIGHT_GRAY)
])


# ════════════════════════════════════════════
# SLIDE 10: 전체 프롬프트 (2/2)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), RED)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
         "📝 재현용 프롬프트 (2/2)", 28, RED, True)

prompt2 = """## 무기 8종
부적(projectile,cd:0.5,레벨=개수,evolved:7발+처치시범위폭발) / 신령방울(homing,cd:0.8,min(lv,3)발)
도깨비방망이(spin,cd:1.0,궤도회전,min(1+lv/2,3)개,evolved:4개+넉백) / 여우불(zone,cd:1.5,evolved:9궤도화염)
천둥(thunder,cd:1.2,즉발min(lv,3)대상,evolved:5체인+슬로우1.5초 게임타임기반)
신궁(pierce,cd:0.6,min(lv,3)발,lv5:무한관통) / 용의숨결(breath,cd:1.8,부채꼴,lv3+화상,evolved:관통드래곤)
귀살검(slash,cd:0.4,전방베기,lv5:후방도공격)

투사체별 개별 크리티컬 판정(rollCrit()), 크리시 "크리!" + 대미지×2

## 패시브 8종: 음양오행(atk+15%), 구미호가죽(spd+10%), 풍백가호(range+12%), 황금(exp+15%),
산삼(regen: 3초/lv간격, 3+lv*2회복, 초록숫자+파티클), 여의주(cdr-8%),
도깨비감투(dodge+10%, "회피!"텍스트+파란파티클), 삼족오부적(crit+8%, "크리!"텍스트)

## 진화: 부적+음양오행=봉인진 / 여우불+구미호가죽=삼매화 / 천둥+풍백가호=뇌신
방망이+황금=여의봉 / 용숨결+여의주=청룡

## 적 8종: 잡귀(hp:2,spd:65,dmg:2,straight,0초) 도깨불(4,50,4,zigzag,30초)
물귀신(6,40,4,aimed,90초) 야차(5,100,5,swooper,150초) 강시(16,30,6,tank,210초)
원귀(6,18,4,sniper투사체,270초) 삼두구(3,55,3,formation3마리,330초) 이무기(12,45,5,spiral,400초)

## 보스: 귀왕(HP:200,SPD:35,DMG:12,450초,추적→돌진→폭발+소환)
구미호왕(HP:500,SPD:30,DMG:15,900초,궤도+유도탄→여우불장판, 처치시 victory화면)

## 밸런스: 스폰=1.2+gameTime/60*0.6+floor(gameTime/60)*0.3
웨이브30초마다=4+floor(gameTime/60), HP스케일=1+gameTime/60*0.1
맵반경1200, 적max200(먼적제거), 폭탄cd30초(전체50대미지,보스25)

## 버그방지: 적uid사용(_tick추적), 천둥슬로우=게임타임타이머(setTimeout금지),
호밍리타겟, 적존대미지=0.5초틱(dmg*dt금지)

game.js가 길면 여러 파트로 나눠서 cat >>으로 이어붙여.
완성 후 GitHub Pages 배포."""

add_multiline(slide, Inches(0.4), Inches(0.9), Inches(12.5), Inches(6.3), [
    (prompt2, 10, LIGHT_GRAY)
])


# ════════════════════════════════════════════
# SLIDE 11: 마무리
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), SW, Inches(0.06), GOLD)

add_multiline(slide, Inches(0), Inches(2.0), SW, Inches(4.0), [
    ("퇴마록 재현 가이드", 36, GOLD, True, PP_ALIGN.CENTER),
    ("", 12),
    ("이 프롬프트로 동일한 게임을 처음부터 재현할 수 있습니다", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 12),
    ("📁  3 파일: index.html + style.css + game.js (~2100줄)", 14, MED_GRAY, False, PP_ALIGN.CENTER),
    ("🎭  6 캐릭터  •  ⚔️ 8 무기  •  💎 8 패시브  •  🔄 5 진화", 14, MED_GRAY, False, PP_ALIGN.CENTER),
    ("👹  8 적 타입  •  💀 2 보스  •  🏆 승리 화면", 14, MED_GRAY, False, PP_ALIGN.CENTER),
    ("", 16),
    ("studiovarem-ui.github.io/toemarok", 14, RGBColor(100, 180, 255), False, PP_ALIGN.CENTER),
])

add_shape(slide, Inches(0), Inches(7.44), SW, Inches(0.06), GOLD)


# ── Save ──
output_path = "/Users/jaehoho/Desktop/j-hoho/ai/game_03/퇴마록_재현_프롬프트.pptx"
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
